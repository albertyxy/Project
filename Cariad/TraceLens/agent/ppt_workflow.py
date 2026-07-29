# -*- coding: utf-8 -*-
"""PPT 报告生成工作流。

确定性流程，无需 LLM。从 Trace Record + XML + MF4 数据生成测试报告 PPT。
作为 agent 模块的一部分，不经过 Planner→Coder 链路。
"""

import os
import re
import copy
import glob
import xml.etree.ElementTree as ET
from typing import Dict, List, Optional, Tuple, Any
from io import BytesIO

import yaml
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image

# 导入 cross_reference 工具
import sys
_tool_scripts_dir = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "Tool Scripts")
)
if _tool_scripts_dir not in sys.path:
    sys.path.insert(0, _tool_scripts_dir)
from cross_reference import cross_reference


# ============================================================================
# 颜色常量
# ============================================================================
COLOR_GREEN = RGBColor(0x00, 0xB0, 0x50)      # Pass
COLOR_ORANGE = RGBColor(0xFF, 0x8C, 0x00)     # Collision
COLOR_RED = RGBColor(0xFF, 0x00, 0x00)        # Failed
COLOR_GRAY = RGBColor(0x99, 0x99, 0x99)       # 无数据

# 场景->Slide 模板映射 (Slide 编号从 1 开始)
SCENARIO_SLIDE_MAP = {
    "CPLA-25_Night_20kph": 2,
    "CPLA-25_Night_40kph": 3,
    "CPLA-25_Night_60kph": 4,
    "CPLA-25_Night_80kph": 5,
    "CPNCO-25_20kph": 6,
    "CPNCO-25_40kph": 7,
    "CBNAO-50_20kph": 8,
    "CBNAO-50_40kph": 9,
    "CBNAO-50_60kph": 10,
    "CPTA-LN-50_10kph": 11,
    "CPTA-LN-50_20kph": 12,
    "CPTA-LN-50_30kph": 13,
    "CPTA-LF-50_10kph": 14,
    "CPTA-LF-50_20kph": 15,
    "CPTA-LF-50_30kph": 16,
    "C2C SCP_30kph": 17,
    "C2C SCP_40kph": 18,
    "C2C SCP_50kph": 19,
    "C2C SCP_60kph": 20,
}

# ECU Name -> XML Systembezeichnung 映射
ECU_NAME_MAP = {
    "LRR": "WBA Bosch PPE",
    "MFK": "MFK5",
    "HCP1": "HCP1 BOSCH EP",
    "NR": "Nanoradar 1",
    "ESC": "ABS",
}


# ============================================================================
# read_trace_record
# ============================================================================
def read_trace_record(excel_path: str) -> Dict[str, Dict[str, Any]]:
    """读取 Trace Record.xlsx，返回场景数据字典。

    Args:
        excel_path: Trace Record.xlsx 文件路径

    Returns:
        {文件夹名: {scenario, vvut, runs: {run1: {trace, result, vimpact}, ...}}}
    """
    wb = openpyxl.load_workbook(excel_path)
    ws = wb["A5L"]

    # 读取所有行，跳过前 3 行（标题行），行 4 开始是表头/数据
    rows = list(ws.iter_rows(min_row=1, values_only=True))

    # 找表头行 (Scenario | Function | VVUT | ...)
    header_row = None
    for i, row in enumerate(rows):
        if row[1] and str(row[1]).strip() == "Scenario":
            header_row = i
            break

    if header_row is None:
        raise ValueError("找不到 Trace Record 表头行")

    data_rows = rows[header_row + 1:]

    # 解析数据：维护当前场景上下文
    current_scenario = None
    current_vvut = None
    current_function = None
    scenario_data: Dict[str, Dict[str, Any]] = {}

    for row in data_rows:
        col_b = str(row[1]).strip() if row[1] else None    # Scenario
        col_c = str(row[2]).strip() if row[2] else None    # Function
        col_d = row[3]                                      # VVUT
        col_g = str(row[6]).strip() if row[6] else None     # Runs (run1/run2/run3)
        col_h = str(row[7]).strip() if row[7] else None     # Pass/No
        col_i = str(row[8]).strip() if row[8] else None     # CANape Trace
        col_l = row[11]                                     # Vimpact

        # 更新当前场景上下文
        if col_b and col_b != "None":
            current_scenario = col_b.replace("\n", "_")
        if col_d is not None and str(col_d) != "None":
            try:
                current_vvut = int(float(col_d))
            except (ValueError, TypeError):
                pass
        if col_c and col_c != "None":
            current_function = col_c

        # 构造文件夹名
        if current_scenario and current_vvut is not None:
            folder_name = f"{current_scenario}_{current_vvut}kph"
        else:
            continue

        # 归一化 Result 值
        result = None
        if col_h and col_h.lower() != "none":
            raw = col_h.strip()
            if raw.lower() == "pass":
                result = "Pass"
            elif raw.lower() == "collision":
                result = "Collision"
            elif raw.lower() == "failed":
                result = "Failed"
            elif raw == "":
                result = None

        # Vimpact 值
        vimpact = None
        if col_l is not None and str(col_l) != "None" and str(col_l).strip() != "":
            try:
                vimpact = float(col_l)
            except (ValueError, TypeError):
                pass

        # Trace 文件名 (去掉 .MF4 后缀如果有)
        trace_name = None
        if col_i and col_i.lower() != "none":
            trace_name = col_i.strip()
            if trace_name.upper().endswith(".MF4"):
                trace_name = trace_name[:-4]

        # 初始化场景条目
        if folder_name not in scenario_data:
            scenario_data[folder_name] = {
                "scenario": current_scenario,
                "vvut": current_vvut,
                "runs": {},
            }

        # 添加 run 数据（只取 run1/run2/run3）
        if col_g and col_g.lower() != "none":
            run_key = col_g.strip().lower()  # run1, run2, run3, run4...
            if run_key in ("run1", "run2", "run3"):
                scenario_data[folder_name]["runs"][run_key] = {
                    "trace": trace_name,
                    "result": result,
                    "vimpact": vimpact,
                }

    wb.close()
    return scenario_data


# ============================================================================
# match_traces
# ============================================================================
def match_traces(
    project_dir: str, trace_data: Dict[str, Dict[str, Any]]
) -> Dict[str, Dict[str, Any]]:
    """匹配场景文件夹与实际 MF4 文件。

    Args:
        project_dir: data/{项目名}/ 目录路径
        trace_data: read_trace_record() 的输出

    Returns:
        {文件夹名: {scenario, vvut, slide_num, runs: {run1: {trace, result,
         vimpact, mf4_path}, ...}}}
        只包含实际存在的场景文件夹
    """
    # 扫描 project_dir 下的 MF4 文件
    mf4_files: Dict[str, str] = {}  # {文件名(无后缀): 完整路径}
    for f in glob.glob(os.path.join(project_dir, "*.MF4")):
        basename = os.path.basename(f)
        name_no_ext = basename[:-4] if basename.upper().endswith(".MF4") else basename
        mf4_files[name_no_ext] = f
    for f in glob.glob(os.path.join(project_dir, "*.mf4")):
        basename = os.path.basename(f)
        name_no_ext = basename[:-4] if basename.lower().endswith(".mf4") else basename
        if name_no_ext not in mf4_files:
            mf4_files[name_no_ext] = f

    # 扫描场景文件夹
    matched: Dict[str, Dict[str, Any]] = {}
    for folder_name, data in trace_data.items():
        scenario_path = os.path.join(project_dir, folder_name)
        if not os.path.isdir(scenario_path):
            continue

        slide_num = SCENARIO_SLIDE_MAP.get(folder_name)
        if slide_num is None:
            continue

        entry = {
            "scenario": data["scenario"],
            "vvut": data["vvut"],
            "slide_num": slide_num,
            "folder_path": scenario_path,
            "runs": {},
        }

        for run_key in ("run1", "run2", "run3"):
            if run_key in data["runs"]:
                run_info = copy.deepcopy(data["runs"][run_key])
                trace_name = run_info.get("trace")
                run_info["mf4_path"] = mf4_files.get(trace_name) if trace_name else None
                entry["runs"][run_key] = run_info

        if entry["runs"]:
            matched[folder_name] = entry

    return matched


# ============================================================================
# Slide 删除工具
# ============================================================================
def _delete_slide_by_index(prs: Presentation, idx: int) -> None:
    """删除指定索引的 slide (0-based)。

    python-pptx 无直接 API，通过操作 XML 实现。
    """
    slide_id = prs.slides[idx].slide_id
    sldId_lst = prs._element.sldIdLst

    rId_to_drop = None
    for sldId_elem in sldId_lst:
        if sldId_elem.get("id") == str(slide_id):
            rId_to_drop = sldId_elem.get(qn("r:id"))
            sldId_lst.remove(sldId_elem)
            break

    if rId_to_drop:
        try:
            prs.part.drop_rel(rId_to_drop)
        except Exception:
            pass


def delete_unused_slides(
    prs: Presentation, active_scenarios: Dict[str, Dict[str, Any]]
) -> None:
    """删除没有对应数据文件夹的场景 slide。

    Args:
        prs: Presentation 对象
        active_scenarios: match_traces() 返回的有效场景字典
    """
    active_slide_nums = {
        data["slide_num"] for data in active_scenarios.values()
    }

    # 收集要删除的 slide 索引 (0-based)
    # Slide 1 (index 0) 由 generate_report 按 XML 有无单独处理
    slides_to_delete = []
    for i in range(1, 20):  # Slide 2-20 -> index 1-19
        slide_num = i + 1
        if slide_num not in active_slide_nums:
            slides_to_delete.append(i)

    for idx in reversed(slides_to_delete):
        _delete_slide_by_index(prs, idx)


# ============================================================================
# populate_slide1
# ============================================================================
def populate_slide1(prs: Presentation, xml_path: str) -> None:
    """填充 Slide 1 的 ECU 配置信息。

    Args:
        prs: Presentation 对象
        xml_path: XML 文件路径
    """
    # 检测 XML 编码（可能是 Windows-1252）
    with open(xml_path, "rb") as f:
        header = f.read(200).decode("utf-8", errors="replace")
    enc_match = re.search(r'encoding="([^"]+)"', header)
    file_encoding = enc_match.group(1) if enc_match else "utf-8"

    with open(xml_path, "r", encoding=file_encoding) as f:
        xml_text = f.read()
    root = ET.fromstring(xml_text)

    # 提取元数据
    vin = ""
    user_projekt = ""

    fahrgestell = root.find(".//Fahrgestellnummer")
    if fahrgestell is not None:
        vin = (fahrgestell.text or "").strip()

    up = root.find(".//UserProjekt")
    if up is not None:
        user_projekt = (up.text or "").strip()

    # 构建 ECU 字典: Systembezeichnung -> {SWTeilenummer, SWVersion, ...}
    # Systembezeichnung 在 Diagnosebloecke/Diagnoseblock 下
    ecu_data: Dict[str, Dict[str, str]] = {}
    for db_elem in root.findall(".//Diagnoseblock"):
        sys_name = db_elem.find("Systembezeichnung")
        if sys_name is None:
            continue
        name = (sys_name.text or "").strip()
        info = {}
        for tag in ("SWTeilenummer", "SWVersion", "HWTeilenummer",
                     "HWVersion", "ZdcName", "ZdcVersion"):
            el = db_elem.find(tag)
            info[tag] = (el.text or "").strip() if el is not None and el.text else ""
        ecu_data[name] = info

    # 填充 Slide 1
    slide1 = prs.slides[0]
    table = None
    for shape in slide1.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        return

    # 更新 r1: Vehicle + VIN (模版样式: 白色粗体, 12pt)
    r1_cell = table.cell(1, 0)
    r1_cell.text = ""
    p = r1_cell.text_frame.paragraphs[0]
    # 空白占位 run (14pt, 模版首个 run 的间距)
    run_spacer = p.add_run()
    run_spacer.text = " "
    run_spacer.font.size = Pt(14)
    run_spacer.font.bold = True
    run_spacer.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # 标签: "Vehicle: ..., VIN: "
    run_label = p.add_run()
    run_label.text = f"Vehicle: {user_projekt}, VIN: "
    run_label.font.size = Pt(12)
    run_label.font.bold = True
    run_label.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # VIN 值
    run_vin = p.add_run()
    run_vin.text = vin
    run_vin.font.size = Pt(12)
    run_vin.font.bold = True
    run_vin.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 更新 ECU 行 r3-r7
    ecu_rows = [
        (3, "LRR"),
        (4, "MFK"),
        (5, "HCP1"),
        (6, "NR"),
        (7, "ESC"),
    ]
    for row_idx, ecu_name in ecu_rows:
        xml_name = ECU_NAME_MAP.get(ecu_name, "")
        info = ecu_data.get(xml_name, {})
        if not info:
            continue

        zdc_name = info.get("ZdcName", "")
        zdc_ver = info.get("ZdcVersion", "")
        zdc_text = f"{zdc_name}[{zdc_ver}]" if zdc_name else ""

        # col 0 (Main-SG) 保持模版样式不动
        # col 1-5: 直接更新现有 run 文本，保留模版字体样式
        _update_cell_text_keep_style(table.cell(row_idx, 1), info.get("SWTeilenummer", ""))
        _update_cell_text_keep_style(table.cell(row_idx, 2), info.get("SWVersion", ""))
        _update_cell_text_keep_style(table.cell(row_idx, 3), zdc_text)
        _update_cell_text_keep_style(table.cell(row_idx, 4), info.get("HWTeilenummer", ""))
        _update_cell_text_keep_style(table.cell(row_idx, 5), info.get("HWVersion", ""))


# ============================================================================
# replace_scenario_images
# ============================================================================
def replace_scenario_images(slide, scenario_folder: str) -> None:
    """替换场景 slide 中的 8 张图片。

    按 shape 位置排序（Y 分行 -> X 排序），1:1 对应 1.PNG..8.PNG。

    Args:
        slide: Slide 对象
        scenario_folder: 场景文件夹路径 (含 PNG 图片)
    """
    # 收集所有图片 shape
    img_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            img_shapes.append(shape)

    if len(img_shapes) != 8:
        return

    # 按 Y 分行 (容差 0.4in = 365760 EMU)
    Y_TOLERANCE_EMU = 365760
    img_shapes.sort(key=lambda s: (s.top, s.left))

    # 分组: 相邻 shape 的 Y 差值 < 容差 → 同一行
    rows = []
    current_row = [img_shapes[0]]
    for s in img_shapes[1:]:
        if s.top - current_row[-1].top < Y_TOLERANCE_EMU:
            current_row.append(s)
        else:
            rows.append(sorted(current_row, key=lambda x: x.left))
            current_row = [s]
    rows.append(sorted(current_row, key=lambda x: x.left))

    # 展开为有序列表
    ordered_shapes = []
    for row in rows:
        ordered_shapes.extend(row)

    # 替换图片
    for i, shape in enumerate(ordered_shapes):
        png_path = os.path.join(scenario_folder, f"{i+1}.PNG")
        if os.path.isfile(png_path):
            # 获取原图片尺寸，保持比例
            img = Image.open(png_path)
            img_w, img_h = img.size
            img.close()
            shape_w = shape.width
            shape_h = shape.height
            # 用 BytesIO 读取
            with open(png_path, "rb") as f:
                new_blob = f.read()
            # 替换图片: 通过 rId 找到 image part 并替换 blob
            pic = shape._element
            blipFill = pic.find(qn('p:blipFill'))
            if blipFill is not None:
                blip = blipFill.find(qn('a:blip'))
                if blip is not None:
                    rId = blip.get(qn('r:embed'))
                    if rId:
                        try:
                            image_part = slide.part.related_part(rId)
                            image_part._blob = new_blob
                        except Exception:
                            pass
        else:
            # PNG 不存在 → 置空
            shape.element.getparent().remove(shape.element)


# ============================================================================
# fill_scenario_table
# ============================================================================
def fill_scenario_table(
    slide,
    scenario_name: str,
    run_data: Dict[str, Any],
    yaml_config: Dict[str, Any],
    project_dir: str,
) -> None:
    """填充场景 slide 的表格。

    Args:
        slide: Slide 对象
        scenario_name: 场景名
        run_data: match_traces() 返回的单个场景数据 (含 runs)
        yaml_config: table_mapping.yaml 加载后的配置
        project_dir: data/{项目名}/ 目录路径
    """
    scenarios_cfg = yaml_config.get("scenarios", {})
    col_cfg = scenarios_cfg.get(scenario_name, {})
    if not col_cfg:
        return

    columns = col_cfg.get("columns", [])
    if not columns:
        return

    # 找到表格 (包括 AlternateContent 中嵌套的表格)
    raw_tbl = _find_table_on_slide(slide)
    if raw_tbl is None:
        return

    runs_data = run_data.get("runs", {})
    n_rows = len(raw_tbl.rows)

    # 分离 cross_reference 列和其他列
    xref_cols = [c for c in columns if c.get("source") in ("", "cross_reference") and c.get("extract")]
    other_cols = [c for c in columns if c not in xref_cols]

    # 处理 fixed / trace_record 列 (快)
    for col_def in other_cols:
        col_idx = col_def["col_index"]
        source = col_def.get("source", "")

        if source == "fixed":
            for run_i, run_key in enumerate(("run1", "run2", "run3")):
                row_idx = run_i + 1
                if row_idx < n_rows:
                    raw_tbl.set_cell_text(row_idx, col_idx, run_key.capitalize())

        elif source == "trace_record":
            for run_i, run_key in enumerate(("run1", "run2", "run3")):
                row_idx = run_i + 1
                if row_idx >= n_rows:
                    continue
                run_info = runs_data.get(run_key, {})

                if col_def.get("header") == "Result":
                    result_val = run_info.get("result")
                    raw_tbl.set_cell_text(row_idx, col_idx,
                                          result_val if result_val else "/")
                    if result_val:
                        raw_tbl.set_cell_color(row_idx, col_idx,
                                               _result_color_hex(result_val))
                elif col_def.get("header", "").startswith("V_impact"):
                    vimpact = run_info.get("vimpact")
                    fmt = col_def.get("format", ".3f")
                    val = f"{vimpact:{fmt}}" if vimpact is not None else "/"
                    raw_tbl.set_cell_text(row_idx, col_idx, val)
                else:
                    raw_tbl.set_cell_text(row_idx, col_idx, "/")

    # 处理 cross_reference 列: 按 run 批量，每 run 只打开一次 MF4
    if xref_cols:
        from asammdf import MDF
        for run_i, run_key in enumerate(("run1", "run2", "run3")):
            row_idx = run_i + 1
            if row_idx >= n_rows:
                continue
            run_info = runs_data.get(run_key, {})
            mf4_path = run_info.get("mf4_path")
            if not mf4_path or not os.path.isfile(mf4_path):
                for col_def in xref_cols:
                    raw_tbl.set_cell_text(row_idx, col_def["col_index"], "/")
                continue

            try:
                with MDF(mf4_path) as mdf:
                    for col_def in xref_cols:
                        extract_cfg = col_def.get("extract", {})
                        target_signal = extract_cfg.get("target", "")
                        triggers = extract_cfg.get("triggers", [])
                        fmt = extract_cfg.get("format", ".2f")
                        col_idx = col_def["col_index"]

                        try:
                            result = cross_reference(
                                file_path=mf4_path,
                                target_signals=[target_signal],
                                triggers=triggers,
                                max_points=1,
                                mdf=mdf,
                            )
                            if result and len(result) > 0:
                                val = result[0]["targets"].get(target_signal)
                                raw_tbl.set_cell_text(row_idx, col_idx,
                                                      f"{float(val):{fmt}}" if val is not None else "/")
                            else:
                                raw_tbl.set_cell_text(row_idx, col_idx, "/")
                        except Exception:
                            raw_tbl.set_cell_text(row_idx, col_idx, "/")
            except Exception:
                for col_def in xref_cols:
                    raw_tbl.set_cell_text(row_idx, col_def["col_index"], "/")


def _update_cell_text_keep_style(cell, text: str) -> None:
    """更新单元格文本，保留模版原有的字体样式 (字号/颜色/粗体等)。"""
    p = cell.text_frame.paragraphs[0]
    runs = list(p.runs)
    if runs:
        first = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
        first.text = text
    else:
        cell.text = text


def _color_first_run(cell, result_val: Optional[str]) -> None:
    """给单元格第一个 run 设置 Result 颜色 (保留模版其他样式)。"""
    runs = list(cell.text_frame.paragraphs[0].runs)
    if not runs:
        return
    if result_val == "Pass":
        runs[0].font.color.rgb = COLOR_GREEN
    elif result_val == "Collision":
        runs[0].font.color.rgb = COLOR_ORANGE
    elif result_val == "Failed":
        runs[0].font.color.rgb = COLOR_RED
    else:
        runs[0].font.color.rgb = COLOR_GRAY


# ============================================================================
# 表格查找 (包括 AlternateContent 中嵌套的表格)
# ============================================================================

# XML 命名空间
_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_AC_NS = 'http://schemas.openxmlformats.org/markup-compatibility/2006'


class _RawTable:
    """直接操作 a:tbl XML 的轻量表格访问器。

    python-pptx 不识别 AlternateContent 中的 graphicFrame 表格，
    此类直接读写 a:tbl XML 元素来更新单元格文本，保留模版样式。
    """

    def __init__(self, tbl_element):
        self._tbl = tbl_element

    @property
    def rows(self):
        return self._tbl.findall(f'{{{_A_NS}}}tr')

    def set_cell_text(self, row_idx: int, col_idx: int, text: str) -> None:
        """设置单元格文本，保留第一个 run 的样式。"""
        tc = self._get_tc(row_idx, col_idx)
        if tc is None:
            return
        txBody = tc.find(f'{{{_A_NS}}}txBody')
        if txBody is None:
            return
        paras = txBody.findall(f'{{{_A_NS}}}p')
        if not paras:
            return
        p = paras[0]
        runs = p.findall(f'{{{_A_NS}}}r')
        if runs:
            first = runs[0]
            for extra in runs[1:]:
                p.remove(extra)
            t_el = first.find(f'{{{_A_NS}}}t')
            if t_el is not None:
                t_el.text = text
        else:
            r_el = p.makeelement(f'{{{_A_NS}}}r', {})
            p.append(r_el)
            t_el = r_el.makeelement(f'{{{_A_NS}}}t', {})
            t_el.text = text
            r_el.append(t_el)

    def set_cell_color(self, row_idx: int, col_idx: int, color_hex: str) -> None:
        """设置单元格第一个 run 的文字颜色 (如 '00B050')。"""
        tc = self._get_tc(row_idx, col_idx)
        if tc is None:
            return
        runs = tc.findall(f'.//{{{_A_NS}}}r')
        if not runs:
            return
        rPr = runs[0].find(f'{{{_A_NS}}}rPr')
        if rPr is None:
            rPr = runs[0].makeelement(f'{{{_A_NS}}}rPr', {})
            runs[0].insert(0, rPr)
        solidFill = rPr.find(f'{{{_A_NS}}}solidFill')
        if solidFill is None:
            solidFill = rPr.makeelement(f'{{{_A_NS}}}solidFill', {})
            rPr.append(solidFill)
        srgbClr = solidFill.find(f'{{{_A_NS}}}srgbClr')
        if srgbClr is None:
            srgbClr = solidFill.makeelement(f'{{{_A_NS}}}srgbClr', {'val': color_hex})
            solidFill.append(srgbClr)
        else:
            srgbClr.set('val', color_hex)

    def _get_tc(self, row_idx: int, col_idx: int):
        tr_list = self._tbl.findall(f'{{{_A_NS}}}tr')
        if row_idx >= len(tr_list):
            return None
        tc_list = tr_list[row_idx].findall(f'{{{_A_NS}}}tc')
        if col_idx >= len(tc_list):
            return None
        return tc_list[col_idx]


def _find_table_on_slide(slide) -> Any:
    """在 slide 上查找表格 (包括 AlternateContent 中嵌套的 graphicFrame)。

    Returns:
        _RawTable 对象，找不到返回 None
    """
    # 路径 1: 普通 shape 表格 (Slide 2-5)
    for shape in slide.shapes:
        if shape.has_table:
            tbl_el = shape._element.find(f'.//{{{_A_NS}}}tbl')
            if tbl_el is not None:
                return _RawTable(tbl_el)

    # 路径 2: AlternateContent 中的 graphicFrame 表格 (Slide 6-20)
    for ac in slide._element.findall(f'.//{{{_AC_NS}}}AlternateContent'):
        for choice in ac:
            for gf in choice.findall(f'.//{{{_P_NS}}}graphicFrame'):
                tbl_el = gf.find(f'.//{{{_A_NS}}}tbl')
                if tbl_el is not None:
                    return _RawTable(tbl_el)

    return None


def _result_color_hex(result_val: Optional[str]) -> str:
    """Result 值对应的 RGB 颜色十六进制字符串。"""
    if result_val == "Pass":
        return "00B050"
    elif result_val == "Collision":
        return "FF8C00"
    elif result_val == "Failed":
        return "FF0000"
    else:
        return "999999"


# ============================================================================
# 文本框格式常量 (匹配模版样式)
# ============================================================================
_TEXT_SIZE = Pt(10)
_TEXT_BODY_COLOR = RGBColor(0x11, 0x11, 0x11)

# 模版文本框段落索引 (固定结构)
_P_SCENARIO_NAME = 1     # Scenario 名
_P_TRACE_RUN1 = 9        # Trace Run1
_P_TRACE_RUN2 = 10       # Trace Run2
_P_TRACE_RUN3 = 11       # Trace Run3
_P_RESULT_RUN1 = 13      # Result Run1
_P_RESULT_RUN2 = 14      # Result Run2
_P_RESULT_RUN3 = 15      # Result Run3
_P_ANALYSIS_TEXT = 17    # Analysis 内容


# ============================================================================
# fill_scenario_text
# ============================================================================
def fill_scenario_text(slide, scenario_name: str, run_data: Dict[str, Any]) -> None:
    """填充场景 slide 的左侧文本框。

    保留模版段落结构，仅更新数据相关段落，匹配模版字体样式。

    模版结构: Scenario → (空行x6) → Trace → Result → Analysis

    Args:
        slide: Slide 对象
        scenario_name: 场景名
        run_data: match_traces() 返回的单个场景数据 (含 runs)
    """
    text_box = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "文本框 2":
            text_box = shape
            break
    if text_box is None:
        return

    tf = text_box.text_frame
    runs_data = run_data.get("runs", {})

    # P1: 场景名
    _replace_para_text(tf.paragraphs[_P_SCENARIO_NAME], scenario_name, bold=False)

    # P9-11: Trace 文件名 (模版: 3-run 结构 label + filename + .MF4)
    for i, run_key in enumerate(("run1", "run2", "run3")):
        para = tf.paragraphs[_P_TRACE_RUN1 + i]
        run_info = runs_data.get(run_key, {})
        trace_name = run_info.get("trace", "/")
        _replace_trace_para(para, run_key, trace_name)

    # P13-15: Result 行 (带颜色)
    for i, run_key in enumerate(("run1", "run2", "run3")):
        para = tf.paragraphs[_P_RESULT_RUN1 + i]
        run_info = runs_data.get(run_key, {})
        _replace_result_para(para, run_key, run_info)

    # P17-P18: Analysis 清空 (有些 slide 跨两段)
    _replace_para_text(tf.paragraphs[_P_ANALYSIS_TEXT], "", bold=False)
    if _P_ANALYSIS_TEXT + 1 < len(tf.paragraphs):
        _replace_para_text(tf.paragraphs[_P_ANALYSIS_TEXT + 1], "", bold=False)


def _clear_para_runs(para):
    """清除段落中所有 run 元素"""
    for run in list(para.runs):
        run._r.getparent().remove(run._r)


def _replace_para_text(para, text: str, bold: bool = False):
    """替换段落文本，保留模版样式。

    Args:
        para: 段落对象
        text: 新文本
        bold: 是否加粗
    """
    _clear_para_runs(para)
    run = para.add_run()
    run.text = text
    run.font.size = _TEXT_SIZE
    run.font.bold = bold
    run.font.color.rgb = _TEXT_BODY_COLOR


def _replace_trace_para(para, run_key: str, trace_name: str):
    """替换 Trace 段落，匹配模版 3-run 结构: label + filename + .MF4。

    模版样式:
      - label (如 "Run1: "): Pt(10), color=111111
      - filename: Pt(10), color=000000
      - ".MF4": Pt(10), color=inherit (无显式颜色)

    Args:
        para: 段落对象
        run_key: run1/run2/run3
        trace_name: MF4 文件名 (无后缀, "/" 表示无数据)
    """
    _clear_para_runs(para)

    if trace_name and trace_name != "/":
        # run 1: 标签 "Run1: "
        r_label = para.add_run()
        r_label.text = f"{run_key.capitalize()}: "
        r_label.font.size = _TEXT_SIZE
        r_label.font.color.rgb = _TEXT_BODY_COLOR

        # run 2: 文件名
        r_name = para.add_run()
        r_name.text = trace_name
        r_name.font.size = _TEXT_SIZE
        r_name.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # run 3: ".MF4" 后缀
        r_ext = para.add_run()
        r_ext.text = ".MF4"
        r_ext.font.size = _TEXT_SIZE
    else:
        # run 1: 标签
        r_label = para.add_run()
        r_label.text = f"{run_key.capitalize()}: "
        r_label.font.size = _TEXT_SIZE
        r_label.font.color.rgb = _TEXT_BODY_COLOR

        # run 2: "/"
        r_slash = para.add_run()
        r_slash.text = "/"
        r_slash.font.size = _TEXT_SIZE
        r_slash.font.color.rgb = COLOR_GRAY


def _replace_result_para(para, run_key: str, run_info: dict):
    """替换 Result 段落，带颜色标记。

    模版 3-run 结构: "Run1" + ":  " + result值(着色)

    Args:
        para: 段落对象
        run_key: run1/run2/run3
        run_info: {trace, result, vimpact, mf4_path}
    """
    _clear_para_runs(para)

    result_val = run_info.get("result", "")
    vimpact = run_info.get("vimpact")

    # run 1: "Run1" (纯标签，不含冒号)
    r_label = para.add_run()
    r_label.text = f"{run_key.capitalize()}"
    r_label.font.size = _TEXT_SIZE
    r_label.font.color.rgb = _TEXT_BODY_COLOR

    # run 2: ":  " (冒号+空格分隔)
    r_sep = para.add_run()
    r_sep.text = ":  "
    r_sep.font.size = _TEXT_SIZE
    r_sep.font.color.rgb = _TEXT_BODY_COLOR

    # run 3: 结果值 (着色)
    if result_val == "Pass":
        r = para.add_run()
        r.text = "Pass"
        r.font.size = _TEXT_SIZE
        r.font.color.rgb = COLOR_GREEN

    elif result_val == "Collision":
        r = para.add_run()
        r.text = "Collision"
        r.font.size = _TEXT_SIZE
        r.font.color.rgb = COLOR_ORANGE

        if vimpact is not None:
            r2 = para.add_run()
            r2.text = f" - V_impact = {vimpact} kph"
            r2.font.size = _TEXT_SIZE
            r2.font.color.rgb = _TEXT_BODY_COLOR

    elif result_val == "Failed":
        r = para.add_run()
        r.text = "Failed"
        r.font.size = _TEXT_SIZE
        r.font.color.rgb = COLOR_RED

    else:
        r = para.add_run()
        r.text = "/"
        r.font.size = _TEXT_SIZE
        r.font.color.rgb = COLOR_GRAY


# ============================================================================
# generate_report
# ============================================================================
def generate_report(
    project_dir: str,
    template_path: str,
    output_path: str,
    yaml_path: Optional[str] = None,
) -> str:
    """主流程：生成 PPT 测试报告。

    Args:
        project_dir: data/{项目名}/ 目录路径
        template_path: PPT 模板文件路径
        output_path: 输出 PPT 文件路径
        yaml_path: table_mapping.yaml 路径 (默认 pptx_agent/table_mapping.yaml)

    Returns:
        输出文件路径
    """
    # 1. 加载 YAML 配置
    if yaml_path is None:
        yaml_path = os.path.join(os.path.dirname(__file__), "table_mapping.yaml")
    with open(yaml_path, "r", encoding="utf-8") as f:
        yaml_config = yaml.safe_load(f)

    # 2. 读取 Trace Record (可选)
    excel_path = os.path.join(project_dir, "Trace Record.xlsx")
    has_trace = os.path.isfile(excel_path)

    matched: Dict[str, Dict[str, Any]] = {}
    if has_trace:
        trace_data = read_trace_record(excel_path)
        print(f"[generate_report] 读取 Trace Record: {len(trace_data)} 个场景")
        # 3. 匹配场景 -> MF4
        matched = match_traces(project_dir, trace_data)
        print(f"[generate_report] 匹配到 {len(matched)} 个有效场景")

    # 4. 打开模板
    prs = Presentation(template_path)

    # 5. 删除无数据的 slide (全部场景 slide 删除时，只保留 Slide 1)
    delete_unused_slides(prs, matched)
    print(f"[generate_report] 剩余 {len(prs.slides)} 页")

    # 6. 查找 XML 文件并填充 Slide 1，无 XML 则删除 Slide 1
    xml_files = glob.glob(os.path.join(project_dir, "*.xml"))
    if xml_files:
        xml_path = xml_files[0]
        populate_slide1(prs, xml_path)
        print(f"[generate_report] Slide 1 ECU 配置已填充")
    else:
        _delete_slide_by_index(prs, 0)
        print(f"[generate_report] 无 XML 文件，Slide 1 已删除")

    # 7. 填充每个场景 slide (仅在 matched 非空时)
    for folder_name, data in matched.items():
        slide_num = data["slide_num"]
        # 删除 slide 后索引可能偏移，需要用 slide_id 定位
        slide = _find_slide_by_number(prs, slide_num)
        if slide is None:
            continue

        scenario_name = folder_name
        folder_path = data["folder_path"]

        # 替换图片
        replace_scenario_images(slide, folder_path)

        # 填文本框
        fill_scenario_text(slide, scenario_name, data)

        # 填表格
        fill_scenario_table(slide, scenario_name, data, yaml_config, project_dir)

        print(f"[generate_report] Slide {slide_num}: {scenario_name} 完成")

    # 8. 保存
    prs.save(output_path)
    print(f"[generate_report] 报告已保存: {output_path}")

    return output_path


def _find_slide_by_number(prs: Presentation, slide_num: int):
    """根据 slide 编号查找 slide 对象 (删除操作后索引可能变化)。"""
    # Slide 编号从 Slide Number Placeholder 的文本判断
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "Slide Number Placeholder 5":
                text = shape.text_frame.text.strip()
                try:
                    if int(text) == slide_num:
                        return slide
                except ValueError:
                    continue
    return None
