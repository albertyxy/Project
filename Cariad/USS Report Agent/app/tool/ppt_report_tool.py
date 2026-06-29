"""
PPT 测试报告生成工具

封装 python-pptx 操作，提供简洁的 Tool 接口给 Agent 调用。

图片替换策略：
  - 非 TS_FCT_USS_18 Slide (3-8): 简单 1:1 替换所有图片
  - TS_FCT_USS_18 Slide (9-17): 三列布局
      * Example 列 (左, ~1.2in) → 保留模板图片不动
      * Test 列 (中, ~3.7in) → 替换为 IMG_XXXX.JPG 照片
      * Test Result 列 (右, ~6.3in) → 替换为 01.PNG 图表
"""

import io
import json
import os
import re
import shutil
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from app.config import PROJECT_ROOT, WORKSPACE_ROOT
from app.tool.base import BaseTool
from app.utils.logger import logger


class FolderImageCursor:
    """管理数据文件夹中图片的消耗指针。

    对 TS_FCT_USS_18 文件夹，JPG 和 PNG 分别维护独立指针。
    对普通文件夹，使用 all_images 统一指针。
    """

    def __init__(self, folder_path: str):
        self.folder_path = folder_path
        # 统一列表（非 TS_FCT_USS_18 使用）
        self.all_images: List[str] = []
        # JPG/PNG 独立通道（TS_FCT_USS_18 使用）
        self.jpg_images: List[str] = []
        self.png_images: List[str] = []
        self.jpg_consumed: int = 0
        self.png_consumed: int = 0
        self.consumed: int = 0

        if os.path.isdir(folder_path):
            files = sorted(os.listdir(folder_path))
            # 分离 JPG 和 PNG
            for f in files:
                full = os.path.join(folder_path, f)
                if f.lower().endswith((".jpg", ".jpeg")):
                    self.jpg_images.append(full)
                    self.all_images.append(full)
                elif f.lower().endswith(".png"):
                    self.png_images.append(full)
                    self.all_images.append(full)

    # -- 统一接口（非 TS_FCT_USS_18） --
    def take(self, count: int) -> List[str]:
        result = self.all_images[self.consumed : self.consumed + count]
        self.consumed += len(result)
        return result

    # -- JPG/PNG 独立接口（TS_FCT_USS_18） --
    def take_jpg(self, count: int) -> List[str]:
        result = self.jpg_images[self.jpg_consumed : self.jpg_consumed + count]
        self.jpg_consumed += len(result)
        return result

    def take_png(self, count: int) -> List[str]:
        result = self.png_images[self.png_consumed : self.png_consumed + count]
        self.png_consumed += len(result)
        return result

    def take_last_jpg(self, count: int) -> List[str]:
        if count <= 0:
            return []
        return self.jpg_images[-count:] if count <= len(self.jpg_images) else self.jpg_images[:]

    def take_last_png(self, count: int) -> List[str]:
        if count <= 0:
            return []
        return self.png_images[-count:] if count <= len(self.png_images) else self.png_images[:]

    @property
    def jpg_remaining(self) -> int:
        return len(self.jpg_images) - self.jpg_consumed

    @property
    def png_remaining(self) -> int:
        return len(self.png_images) - self.png_consumed


class PPTReportTool(BaseTool):
    """PPT 测试报告生成工具。"""

    name: str = "ppt_report"
    description: str = (
        "PPT测试报告生成工具。可以分析模板PPT、自动匹配测试数据、"
        "替换Slide中的图片、生成最终报告。支持的操作："
        "find_project(根据项目名搜索数据目录和模板)、"
        "analyze_template(分析PPT结构)、"
        "build_mapping(建立Slide与数据的映射)、"
        "generate_report(一键生成完整报告)"
    )

    parameters: dict = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": [
                    "find_project",
                    "analyze_template",
                    "build_mapping",
                    "generate_report",
                    "get_slide_info",
                ],
                "description": "要执行的操作",
            },
            "project_name": {
                "type": "string",
                "description": "项目名称（用于 find_project 操作）",
            },
            "data_root": {
                "type": "string",
                "description": "数据根目录（可选，默认 USS Data；用于 find_project 指定其他模板目录）",
            },
            "template_path": {
                "type": "string",
                "description": "模板PPT文件路径",
            },
            "data_dir": {
                "type": "string",
                "description": "测试数据根目录路径",
            },
            "output_path": {
                "type": "string",
                "description": "输出PPT文件路径（可选，默认保存到workspace/）",
            },
        },
        "required": ["action"],
    }

    # 路径均相对于项目根目录，不依赖 cwd
    @property
    def uss_data_root(self) -> Path:
        return PROJECT_ROOT / "USS Data"

    def _resolve_data_root(self, data_root: Optional[str] = None) -> Path:
        """解析数据根目录。传入路径则相对于 PROJECT_ROOT，默认 USS Data。"""
        if data_root:
            p = Path(data_root)
            return p if p.is_absolute() else PROJECT_ROOT / p
        return self.uss_data_root

    @property
    def uss_template_path(self) -> Path:
        """默认模板路径（USS 报告），新模板在 find_project 中按 data_root 返回。"""
        data_root = self._resolve_data_root()
        cfg = self._load_template_config(data_root)
        tpl = cfg.get("template_file", "USSDB Test Report Template.pptx")
        return data_root / tpl

    # ---- 模板配置 ----
    _template_config: Optional[dict] = None
    _template_config_root: Optional[Path] = None

    def _load_template_config(self, data_root: Optional[Path] = None) -> dict:
        """加载 {data_root}/template_config.json，按 data_root 缓存。"""
        if data_root is None:
            data_root = self._resolve_data_root()
        if self._template_config is not None and self._template_config_root == data_root:
            return self._template_config
        cfg_path = data_root / "template_config.json"
        try:
            if cfg_path.exists():
                with open(cfg_path, "r", encoding="utf-8") as f:
                    self._template_config = json.load(f)
                self._template_config_root = data_root
                logger.info(f"已加载模板配置: {cfg_path}")
            else:
                self._template_config = {}
                self._template_config_root = data_root
        except Exception as e:
            logger.warning(f"模板配置加载失败: {e}")
            self._template_config = {}
            self._template_config_root = data_root
        return self._template_config

    def _cfg(self, *keys, default=None):
        """安全读取嵌套配置。_cfg('slide_layouts', 'default', 'strategy')"""
        cfg = self._load_template_config()
        for k in keys:
            if isinstance(cfg, dict):
                cfg = cfg.get(k)
            else:
                return default
        return cfg if cfg is not None else default

    @property
    def driver_activity_re(self) -> str:
        """Driver Activity 正则：优先从 config 读取，降级使用默认值。"""
        return self._cfg("driver_activity", "regex",
                         default=r"Driver Activity:\s*TS_FCT_USS_\d+(?:\s*[-&]?\s*Action\s*-?\d+)*(?:\s*&\s*TS_FCT_USS_\d+(?:\s*[-&]?\s*Action\s*-?\d+)*)*")

    async def execute(self, **kwargs) -> Any:
        action = kwargs.get("action", "")
        if action == "find_project":
            result = self._find_project(
                kwargs.get("project_name", ""),
                kwargs.get("data_root"),
            )
        elif action == "analyze_template":
            result = self._analyze_template(kwargs.get("template_path", ""))
        elif action == "build_mapping":
            result = self._build_mapping(
                kwargs.get("template_path", ""), kwargs.get("data_dir", "")
            )
        elif action == "generate_report":
            result = self._generate_report(
                kwargs.get("template_path", ""),
                kwargs.get("data_dir", ""),
                kwargs.get("output_path"),
            )
        elif action == "get_slide_info":
            result = self._get_slide_info(
                kwargs.get("template_path", ""), kwargs.get("data_dir", "")
            )
        else:
            return self.fail_response(f"未知操作: {action}")
        return self.success_response(result)

    # ==================================================================
    # 核心方法
    # ==================================================================

    def _find_project(self, project_name: str, data_root: Optional[str] = None) -> dict:
        """在 data_root 下搜索项目文件夹。默认 USS Data，可通过参数指定其他根目录。"""
        root = self._resolve_data_root(data_root)
        if not root.exists():
            return {"found": False, "error": f"数据根目录不存在: {root}"}

        subdirs = [d for d in root.iterdir() if d.is_dir()]
        subdir_names = [d.name for d in subdirs]

        if not project_name.strip():
            return {"found": False, "error": "请提供项目名称", "available_projects": subdir_names}

        normalized_input = self._normalize_name(project_name)
        matches = []
        for d in subdirs:
            if normalized_input in self._normalize_name(d.name) or self._normalize_name(d.name) in normalized_input:
                matches.append(d)

        if len(matches) == 0:
            return {"found": False, "error": f"未找到匹配的项目: '{project_name}'", "available_projects": subdir_names}
        if len(matches) > 1:
            return {"found": False, "error": f"找到多个匹配: {[m.name for m in matches]}，请更具体"}

        project_dir = matches[0]
        data_folders = sorted([d.name for d in project_dir.iterdir() if d.is_dir()])

        # 模板路径：从 data_root 的 config 读取（不限于 USS Data）
        cfg = self._load_template_config(root)
        tpl_name = cfg.get("template_file", "USSDB Test Report Template.pptx")
        tpl_path = root / tpl_name
        if not tpl_path.exists():
            return {"found": False, "error": f"模板文件不存在: {tpl_path}"}

        return {
            "found": True,
            "project_name": project_dir.name,
            "project_dir": str(project_dir),
            "template_path": str(tpl_path),
            "data_root": str(root),
            "data_folders": data_folders,
        }

    def _analyze_template(self, template_path: str) -> dict:
        if not template_path or not Path(template_path).exists():
            return {"error": f"模板文件不存在: {template_path}"}

        prs = Presentation(template_path)
        slides_info = []
        total_images = 0
        total_test_slides = 0

        for i, slide in enumerate(prs.slides, 1):
            image_count = 0
            images_detail = []
            texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        texts.append(txt)
                if shape.shape_type == 13:  # Picture
                    image_count += 1
                    images_detail.append({
                        "index": image_count,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "name": shape.name,
                    })

            total_images += image_count

            driver_activities = []
            for text in texts:
                da_match = re.search(self.driver_activity_re, text, re.IGNORECASE)
                if da_match:
                    for act in self._parse_driver_activity_line(da_match.group(0)):
                        if act not in driver_activities:
                            driver_activities.append(act)

            if driver_activities:
                total_test_slides += 1

            # 检测是否为 TS_FCT_USS_18 系列 Slide（三列布局）
            is_ts_18 = any("TS_FCT_USS_18" in a for a in driver_activities)

            slides_info.append({
                "slide_num": i,
                "image_count": image_count,
                "images": images_detail,
                "texts": texts,
                "driver_activities": driver_activities,
                "is_test_slide": len(driver_activities) > 0,
                "is_ts_18": is_ts_18,
            })

        return {
            "total_slides": len(prs.slides),
            "total_images": total_images,
            "total_test_slides": total_test_slides,
            "slides": slides_info,
        }

    def _build_mapping(self, template_path: str, data_dir: str) -> dict:
        if not template_path or not Path(template_path).exists():
            return {"error": f"模板文件不存在: {template_path}"}
        if not data_dir or not Path(data_dir).exists():
            return {"error": f"数据目录不存在: {data_dir}"}

        template_info = self._analyze_template(template_path)
        if "error" in template_info:
            return template_info

        data_folders = {d.name: str(d) for d in Path(data_dir).iterdir() if d.is_dir()}

        mappings = []
        non_test_count = 0

        for slide_info in template_info["slides"]:
            activities = slide_info["driver_activities"]
            is_ts_18 = slide_info.get("is_ts_18", False)

            if not activities:
                non_test_count += 1
                mappings.append({
                    "slide_num": slide_info["slide_num"],
                    "mapping_type": "non_test",
                    "driver_activities": [],
                    "data_folder": None,
                    "column_mode": False,
                })
                continue

            if len(activities) == 1:
                data_folder_name = activities[0]
                folder_path = data_folders.get(data_folder_name)
                if not folder_path:
                    matched = self._fuzzy_match_folder(data_folder_name, list(data_folders.keys()))
                    if matched:
                        folder_path = data_folders[matched]
                        data_folder_name = matched

                if folder_path:
                    mappings.append({
                        "slide_num": slide_info["slide_num"],
                        "mapping_type": "A/B",
                        "driver_activities": activities,
                        "data_folder": data_folder_name,
                        "data_folder_path": folder_path,
                        "template_image_count": slide_info["image_count"],
                        "column_mode": is_ts_18,
                    })
                else:
                    mappings.append({
                        "slide_num": slide_info["slide_num"],
                        "mapping_type": "unmatched",
                        "driver_activities": activities,
                        "data_folder": None,
                        "column_mode": is_ts_18,
                        "error": f"未找到匹配的数据文件夹: {data_folder_name}",
                    })
            else:
                # 跨文件夹组合（类型 C）
                matched_folders = []
                for act in activities:
                    if act in data_folders:
                        matched_folders.append(act)
                    else:
                        matched = self._fuzzy_match_folder(act, list(data_folders.keys()))
                        matched_folders.append(matched if matched else act)

                mappings.append({
                    "slide_num": slide_info["slide_num"],
                    "mapping_type": "C",
                    "driver_activities": activities,
                    "data_folders": matched_folders,
                    "data_folder_paths": [data_folders.get(f, "") for f in matched_folders],
                    "template_image_count": slide_info["image_count"],
                    "column_mode": is_ts_18,
                })

        # 重新分类 A/B
        self._classify_ab_mappings(mappings)
        type_a_count = sum(1 for m in mappings if m["mapping_type"] == "A")
        type_b_count = sum(1 for m in mappings if m["mapping_type"] == "B")
        type_c_count = sum(1 for m in mappings if m["mapping_type"] == "C")
        type_b_groups = self._find_shared_groups(mappings)

        # 统计列模式数量
        column_mode_count = sum(1 for m in mappings if m.get("column_mode"))
        simple_count = sum(1 for m in mappings if m["mapping_type"] in ("A", "B", "C") and not m.get("column_mode"))

        return {
            "mapping_type_counts": {
                "type_a": type_a_count,
                "type_b": type_b_count,
                "type_b_groups": type_b_groups,
                "type_c": type_c_count,
                "non_test": non_test_count,
                "column_mode": column_mode_count,
                "simple_mode": simple_count,
            },
            "mappings": mappings,
            "data_folders_available": sorted(data_folders.keys()),
        }

    def _generate_report(
        self, template_path: str, data_dir: str, output_path: Optional[str] = None
    ) -> dict:
        if not template_path or not Path(template_path).exists():
            return {"error": f"模板文件不存在: {template_path}"}
        if not data_dir or not Path(data_dir).exists():
            return {"error": f"数据目录不存在: {data_dir}"}

        mapping_result = self._build_mapping(template_path, data_dir)
        if "error" in mapping_result:
            return mapping_result

        mappings = mapping_result["mappings"]

        # 构建 cursor 池
        cursors: Dict[str, FolderImageCursor] = {}
        for d in Path(data_dir).iterdir():
            if d.is_dir():
                cursors[d.name] = FolderImageCursor(str(d))

        # 确定输出路径
        if not output_path:
            project_name = Path(data_dir).name
            WORKSPACE_ROOT.mkdir(parents=True, exist_ok=True)
            output_path = str(WORKSPACE_ROOT / f"{project_name}_Report.pptx")

        shutil.copy2(template_path, output_path)
        prs = Presentation(output_path)

        # Step 0: 自动解析 XML 并填充 Slide 1/2 元数据
        metadata_result = self._populate_metadata(prs, data_dir)

        stats = {
            "metadata": metadata_result,
            "simple_replaced": 0,
            "column_replaced": 0,
            "column_test_replaced": 0,
            "column_test_result_replaced": 0,
            "total_images_replaced": 0,
            "errors": [],
        }

        # 按 Slide 顺序预计算所有列模式图片分配（确保 Type C 在 Type B 之前消耗）
        all_column_allocations = self._precompute_all_columns(mappings, cursors, prs)

        for mapping in mappings:
            slide_num = mapping["slide_num"]
            mtype = mapping["mapping_type"]
            if mtype in ("non_test", "unmatched"):
                if mtype == "unmatched":
                    stats["errors"].append(f"Slide {slide_num}: {mapping.get('error', '未匹配')}")
                continue

            slide = prs.slides[slide_num - 1]
            column_mode = mapping.get("column_mode", False)

            if column_mode:
                alloc = all_column_allocations.get(slide_num, {})
                jpg_list = alloc.get("jpg", [])
                png_list = alloc.get("png", [])
                replaced = self._apply_column_replacement(slide, jpg_list, png_list)
                stats["column_replaced"] += replaced["total"]
                stats["column_test_replaced"] += replaced["test"]
                stats["column_test_result_replaced"] += replaced["test_result"]
                stats["total_images_replaced"] += replaced["total"]
            else:
                replaced = self._replace_simple_mode(slide, mapping, cursors)
                stats["simple_replaced"] += replaced
                stats["total_images_replaced"] += replaced

        prs.save(output_path)

        return {
            "success": True,
            "output_path": output_path,
            "statistics": stats,
            "mapping_summary": mapping_result["mapping_type_counts"],
        }

    def _get_slide_info(self, template_path: str, slide_num: int) -> dict:
        if not template_path or not Path(template_path).exists():
            return {"error": f"模板文件不存在: {template_path}"}
        prs = Presentation(template_path)
        if slide_num < 1 or slide_num > len(prs.slides):
            return {"error": f"Slide 编号超出范围: {slide_num}"}
        slide = prs.slides[slide_num - 1]
        info = {"slide_num": slide_num, "shapes_count": len(slide.shapes), "shapes": []}
        for shape in slide.shapes:
            si = {"name": shape.name, "shape_type": str(shape.shape_type)}
            if shape.has_text_frame:
                si["text"] = shape.text_frame.text.strip()
            info["shapes"].append(si)
        return info

    # ==================================================================
    # XML 元数据解析与填充（Slide 1 标题 + Slide 2 表格）
    # ==================================================================

    def _find_xml_file(self, data_dir: str) -> Optional[str]:
        """在项目数据目录下搜索 XML 文件（支持配置指定 pattern）。"""
        data_path = Path(data_dir)
        if not data_path.exists():
            return None
        pattern = self._cfg("data_sources", "xml", "file_pattern", default="*.xml")
        xml_files = list(data_path.glob(pattern))
        if not xml_files:
            xml_files = list(data_path.glob("*.xml")) + list(data_path.glob("*.XML"))
        return str(xml_files[0]) if xml_files else None

    def _parse_xml_metadata(self, xml_path: str) -> dict:
        """解析 iDEX XML 文件，按配置提取 ECU 元数据。"""
        result = {
            "vehicle_model": "",
            "vin": "",
            "ecus": {},      # {key: {ecu_data}} 如 {"hcp1": {...}, "eps": {...}}
            "xml_file": os.path.basename(xml_path),
        }

        try:
            tree = ET.parse(xml_path)
            root = tree.getroot()

            user_projekt = self._xml_text(root, "UserProjekt")
            result["vehicle_model"] = self._extract_vehicle_model(user_projekt)
            result["vin"] = self._xml_text(root, "Fahrgestellnummer")

            # 从 config 读取 ECU 匹配规则
            ecu_rules = {}
            config_rows = self._cfg("config_slide", "rows", default={})
            for key, row_cfg in config_rows.items():
                match = row_cfg.get("ecu_match")
                if match:
                    ecu_rules[key] = match

            # 遍历 Diagnoseblock 按规则匹配
            for db in root.iter("Diagnoseblock"):
                sys_name = self._xml_text(db, "Systembezeichnung")
                if not sys_name:
                    continue
                for ecu_key, rule in ecu_rules.items():
                    field = rule.get("field", "Systembezeichnung")
                    pattern = rule.get("pattern", "")
                    val = self._xml_text(db, field) if field == "Systembezeichnung" else self._xml_text(db, field)
                    if pattern and pattern.upper() in val.upper():
                        result["ecus"][ecu_key] = self._extract_ecu_data(db)

        except Exception as e:
            logger.warning(f"XML 解析失败: {e}")
            result["error"] = str(e)

        return result

    def _populate_metadata(self, prs, data_dir: str) -> dict:
        """按 template_config.json 配置填充 Slide 1/2 元数据。"""
        result = {"xml_found": False, "title_updated": False, "table_updated": False}

        xml_path = self._find_xml_file(data_dir)
        if not xml_path:
            return result

        result["xml_found"] = True
        meta = self._parse_xml_metadata(xml_path)
        if meta.get("error"):
            result["error"] = meta["error"]
            return result

        # -- Slide 1: 替换标题（配置驱动） --
        title_cfg = self._cfg("title_slide", default={})
        title_slide_idx = title_cfg.get("slide_index", 1) - 1
        if title_slide_idx < len(prs.slides):
            slide1 = prs.slides[title_slide_idx]
            model = meta.get("vehicle_model", "")
            if model:
                for shape in slide1.shapes:
                    if shape.has_text_frame:
                        full_text = shape.text_frame.text
                        if "USSDB" in full_text and "Test Report" in full_text:
                            para = shape.text_frame.paragraphs[0]
                            ussdb_run_idx = None
                            for ri, run in enumerate(para.runs):
                                if "USSDB" in run.text:
                                    ussdb_run_idx = ri
                                    break
                            if ussdb_run_idx is not None and ussdb_run_idx > 0:
                                para.runs[0].text = f"{model} "
                                for ri in range(1, ussdb_run_idx):
                                    para.runs[ri].text = ""
                                result["title_updated"] = True
                            break

        # -- Slide 2: 替换表格数据（配置驱动） --
        config_cfg = self._cfg("config_slide", default={})
        config_slide_idx = config_cfg.get("slide_index", 2) - 1
        if config_slide_idx < len(prs.slides):
            slide2 = prs.slides[config_slide_idx]
            for shape in slide2.shapes:
                if not shape.has_table:
                    continue
                table = shape.table

                config_rows = config_cfg.get("rows", {})
                for key, row_cfg in config_rows.items():
                    ri = row_cfg.get("row_index", -1)
                    if ri < 0 or ri >= len(table.rows):
                        continue

                    if "template" in row_cfg:
                        # 模板字符串行（如 Vehicle + VIN）
                        tmpl = row_cfg["template"]
                        ci = row_cfg.get("col_index", 0)
                        cell = table.cell(ri, ci)
                        text = tmpl.format(
                            vehicle_model=meta.get("vehicle_model", ""),
                            vin=meta.get("vin", ""),
                        )
                        tf = cell.text_frame
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            tf.paragraphs[0].runs[0].text = text
                            for run in tf.paragraphs[0].runs[1:]:
                                run.text = ""
                        else:
                            cell.text = text
                        result["table_updated"] = True

                    elif "ecu_match" in row_cfg:
                        # ECU 数据行
                        ecu_data = meta.get("ecus", {}).get(key, {})
                        if ecu_data:
                            col_map = {}
                            for ci_str, col_cfg in row_cfg.get("columns", {}).items():
                                fmt = col_cfg.get("format", "{value}")
                                val = self._strip(ecu_data.get(col_cfg.get("xml_field", ""), ""))
                                suffix = ""
                                if "suffix_field" in col_cfg:
                                    suffix = self._strip(ecu_data.get(col_cfg["suffix_field"], ""))
                                col_map[int(ci_str)] = fmt.format(value=val, suffix=suffix)
                            self._fill_ecu_row_from_map(table, ri, col_map)
                            result["table_updated"] = True

        return result

    def _fill_ecu_row_from_map(self, table, row_idx: int, col_map: Dict[int, str]):
        """按 {col_index: value} 填入表格行，保持字体样式。"""
        if row_idx >= len(table.rows):
            return
        cols = len(table.columns)
        for ci, val in col_map.items():
            if ci < cols:
                cell = table.cell(row_idx, ci)
                tf = cell.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = val
                    for run in tf.paragraphs[0].runs[1:]:
                        run.text = ""
                else:
                    cell.text = val

    def _fill_ecu_row(self, table, row_idx: int, ecu_data: dict):
        """将 ECU 数据填入表格的指定行，保持原 cell 的字体样式。"""
        if row_idx >= len(table.rows):
            return
        cols = len(table.columns)

        mappings = {
            1: self._strip(ecu_data.get("SWTeilenummer", "")),
            2: self._strip(ecu_data.get("SWVersion", "")),
            3: f"{self._strip(ecu_data.get('ZdcName', ''))} [{self._strip(ecu_data.get('ZdcVersion', ''))}]",
            4: self._strip(ecu_data.get("HWTeilenummer", "")),
            5: self._strip(ecu_data.get("HWVersion", "")),
        }

        for ci, val in mappings.items():
            if ci < cols:
                cell = table.cell(row_idx, ci)
                tf = cell.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = val
                    # 清除后续 runs
                    for run in tf.paragraphs[0].runs[1:]:
                        run.text = ""
                else:
                    cell.text = val

    def _extract_ecu_data(self, db) -> dict:
        """从 Diagnoseblock XML 元素中提取 ECU 数据。"""
        return {
            "Systembezeichnung": self._xml_text(db, "Systembezeichnung"),
            "SWTeilenummer": self._xml_text(db, "SWTeilenummer"),
            "SWVersion": self._xml_text(db, "SWVersion"),
            "HWTeilenummer": self._xml_text(db, "HWTeilenummer"),
            "HWVersion": self._xml_text(db, "HWVersion"),
            "ZdcName": self._xml_text(db, "ZdcName"),
            "ZdcVersion": self._xml_text(db, "ZdcVersion"),
        }

    def _xml_text(self, parent, tag: str) -> str:
        """安全获取 XML 子元素文本。"""
        child = parent.find(tag)
        return (child.text or "").strip() if child is not None and child.text else ""

    def _strip(self, val: str) -> str:
        """去除字符串首尾空白。"""
        return val.strip() if val else ""

    def _extract_vehicle_model(self, user_projekt: str) -> str:
        """从 UserProjekt 提取车型名，优先使用 config 中的正则规则。"""
        if not user_projekt:
            return ""
        rule = self._cfg("title_slide", "fields", "vehicle_model", "transform_rule",
                         default=r"\b(E\d+)\s*(L)?")
        m = re.search(rule, user_projekt, re.IGNORECASE)
        if m:
            base = m.group(1)
            suffix = m.group(2) or ""
            return f"{base}{suffix}"
        return user_projekt

    # ==================================================================
    # 简单模式替换（非 TS_FCT_USS_18）：逐个替换全部图片
    # ==================================================================

    def _replace_simple_mode(self, slide, mapping: dict, cursors: Dict[str, FolderImageCursor]) -> int:
        folder_name = mapping.get("data_folder")
        if not folder_name or folder_name not in cursors:
            return 0

        cursor = cursors[folder_name]
        images_to_use = cursor.take(mapping.get("template_image_count", 0))
        return self._replace_images_1to1(slide, images_to_use)

    # ==================================================================
    # 列模式替换（TS_FCT_USS_18）：Example 保留，Test→JPG，Test Result→PNG
    # ==================================================================

    def _apply_column_replacement(self, slide, jpg_list: List[str], png_list: List[str]) -> dict:
        """执行列模式替换：Example 跳过，Test→JPG，Test Result→PNG。"""
        col_assignments = self._classify_image_columns(slide)
        test_replaced = 0
        tr_replaced = 0

        for i, (pic, col_type) in enumerate(zip(col_assignments["shapes"], col_assignments["column_types"])):
            if col_type == "example":
                continue

            if col_type == "test":
                jpg_idx = sum(1 for j in range(i) if col_assignments["column_types"][j] == "test")
                if jpg_idx < len(jpg_list):
                    img_path = jpg_list[jpg_idx]
                    if os.path.exists(img_path):
                        self._replace_single_image(slide, pic, img_path)
                        test_replaced += 1

            elif col_type == "test_result":
                png_idx = sum(1 for j in range(i) if col_assignments["column_types"][j] == "test_result")
                if png_idx < len(png_list):
                    img_path = png_list[png_idx]
                    if os.path.exists(img_path):
                        self._replace_single_image(slide, pic, img_path)
                        tr_replaced += 1

        return {"total": test_replaced + tr_replaced, "test": test_replaced, "test_result": tr_replaced}

    # ==================================================================
    # 列分类：将 slide 中图片按 left 位置分为 Example / Test / Test Result
    # ==================================================================

    def _classify_image_columns(self, slide) -> dict:
        """按 left 坐标聚类，返回每张图片的列类型。

        Returns:
            {"shapes": [shape, ...], "column_types": ["example"|"test"|"test_result", ...]}
        """
        pic_shapes = []
        lefts = []
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pic_shapes.append(shape)
                lefts.append(shape.left)

        if not pic_shapes:
            return {"shapes": [], "column_types": []}

        # 按 left 坐标排序
        sorted_indices = sorted(range(len(lefts)), key=lambda i: lefts[i])
        sorted_shapes = [pic_shapes[i] for i in sorted_indices]
        sorted_lefts = [lefts[i] for i in sorted_indices]

        # 聚类：相邻 left 值在容差内归为一组
        clusters = []
        current_cluster = [0]
        for i in range(1, len(sorted_lefts)):
            tolerance = self._cfg("slide_layouts", "ts_fct_uss_18", "column_clustering", "tolerance_emu", default=150000)
            if sorted_lefts[i] - sorted_lefts[i - 1] <= tolerance:
                current_cluster.append(i)
            else:
                clusters.append(current_cluster)
                current_cluster = [i]
        clusters.append(current_cluster)

        # 按 left 均值排序 clusters
        cluster_info = []
        for cl in clusters:
            avg_left = sum(sorted_lefts[j] for j in cl) / len(cl)
            cluster_info.append({"indices": cl, "avg_left": avg_left})
        cluster_info.sort(key=lambda c: c["avg_left"])

        # 分配列类型
        num_clusters = len(cluster_info)
        col_names = {}
        if num_clusters == 3:
            col_names = {0: "example", 1: "test", 2: "test_result"}
        elif num_clusters == 2:
            col_names = {0: "test", 1: "test_result"}
        else:
            # 全部归类为 test_result
            for ci in range(num_clusters):
                col_names[ci] = "test_result"

        # 构建结果
        column_types = [""] * len(pic_shapes)
        for ci, cinfo in enumerate(cluster_info):
            ctype = col_names.get(ci, "test_result")
            for idx in cinfo["indices"]:
                column_types[idx] = ctype

        return {"shapes": pic_shapes, "column_types": column_types}

    def _count_column_images(self, slide) -> dict:
        """统计 Slide 中各列的图片数量。"""
        result = self._classify_image_columns(slide)
        counts = {"example": 0, "test": 0, "test_result": 0}
        for ct in result["column_types"]:
            if ct in counts:
                counts[ct] += 1
        return counts

    # ==================================================================
    # 预计算列模式分配（类型 B 共享组）
    # ==================================================================

    def _precompute_all_columns(
        self, mappings: List[dict], cursors: Dict[str, FolderImageCursor], prs
    ) -> Dict[int, dict]:
        """按 Slide 顺序逐页预计算所有列模式的图片分配。

        关键：必须严格按 Slide 编号顺序处理，确保：
        - Type C 组合页（Slide 12, 14）在 Type B 共享页之前消耗对应文件夹的头部图片
        - Type B 共享组之后，同一文件夹的指针已正确推进

        Returns: {slide_num: {"jpg": [...], "png": [...]}, ...}
        """
        allocations: Dict[int, dict] = {}

        # 按 slide_num 排序的列模式映射
        column_mappings = sorted(
            [m for m in mappings if m.get("column_mode")],
            key=lambda m: m["slide_num"],
        )

        for mapping in column_mappings:
            slide_num = mapping["slide_num"]
            mtype = mapping["mapping_type"]
            slide = prs.slides[slide_num - 1]
            col_counts = self._count_column_images(slide)
            needed_test = col_counts["test"]
            needed_tr = col_counts["test_result"]

            if mtype == "C":
                # 从两个文件夹取尾部+头部
                folder_names = mapping.get("data_folders", [])
                if len(folder_names) == 2:
                    ca = cursors.get(folder_names[0])
                    cb = cursors.get(folder_names[1])
                    if ca and cb:
                        # JPG: A 尾部 + B 头部（按比例 1:2）
                        a_jpg_count = needed_test // 3
                        b_jpg_count = needed_test - a_jpg_count
                        jpg_list = ca.take_last_jpg(a_jpg_count) + cb.take_jpg(b_jpg_count)

                        # PNG: A 尾部 + B 头部
                        a_png_count = needed_tr // 3
                        b_png_count = needed_tr - a_png_count
                        png_list = ca.take_last_png(a_png_count) + cb.take_png(b_png_count)

                        allocations[slide_num] = {"jpg": jpg_list, "png": png_list}
                    else:
                        allocations[slide_num] = {"jpg": [], "png": []}
                else:
                    allocations[slide_num] = {"jpg": [], "png": []}

            elif mtype == "A":
                # 类型 A 列模式：单文件夹
                folder_name = mapping.get("data_folder")
                cursor = cursors.get(folder_name) if folder_name else None
                if cursor:
                    allocations[slide_num] = {
                        "jpg": cursor.take_jpg(needed_test),
                        "png": cursor.take_png(needed_tr),
                    }
                else:
                    allocations[slide_num] = {"jpg": [], "png": []}

            elif mtype == "B":
                # 类型 B 共享组：从共享文件夹顺序取
                folder_name = mapping.get("data_folder")
                cursor = cursors.get(folder_name) if folder_name else None
                if cursor:
                    allocations[slide_num] = {
                        "jpg": cursor.take_jpg(needed_test),
                        "png": cursor.take_png(needed_tr),
                    }
                else:
                    allocations[slide_num] = {"jpg": [], "png": []}

        return allocations

    # ==================================================================
    # 底层图片替换
    # ==================================================================

    def _replace_single_image(self, slide, old_pic, img_path: str):
        """替换单张图片，保持位置和尺寸。"""
        try:
            actual_path = self._convert_image_if_needed(img_path)
            left, top, width, height = old_pic.left, old_pic.top, old_pic.width, old_pic.height

            old_sp = old_pic._element
            spTree = old_sp.getparent()
            idx = list(spTree).index(old_sp)

            new_pic = slide.shapes.add_picture(actual_path, left, top, width, height)
            new_sp = new_pic._element

            spTree.remove(new_sp)
            spTree.insert(idx, new_sp)
            spTree.remove(old_sp)
        except Exception as e:
            logger.warning(f"替换图片失败: {e}")

    def _replace_images_1to1(self, slide, image_paths: List[str]) -> int:
        """简单模式：按序一一替换 Slide 中所有图片。"""
        pic_shapes = [s for s in slide.shapes if s.shape_type == 13]
        if not pic_shapes or not image_paths:
            return 0

        replaced = 0
        for i, pic in enumerate(pic_shapes):
            if i >= len(image_paths):
                break
            if os.path.exists(image_paths[i]):
                self._replace_single_image(slide, pic, image_paths[i])
                replaced += 1
        return replaced

    # ==================================================================
    # 辅助方法
    # ==================================================================

    def _classify_ab_mappings(self, mappings: List[dict]):
        ab_indices = [i for i, m in enumerate(mappings) if m["mapping_type"] == "A/B"]
        folder_groups: Dict[str, List[int]] = {}
        for idx in ab_indices:
            folder = mappings[idx].get("data_folder")
            if folder:
                folder_groups.setdefault(folder, []).append(idx)

        for folder, indices in folder_groups.items():
            if len(indices) > 1:
                for i, idx in enumerate(indices):
                    mappings[idx]["mapping_type"] = "B"
                    mappings[idx]["shared_group"] = folder
                    mappings[idx]["index_in_group"] = i
                mappings[indices[0]]["group_first"] = True
            else:
                mappings[indices[0]]["mapping_type"] = "A"
                mappings[indices[0]].pop("data_image_count", None)

    def _find_shared_groups(self, mappings: List[dict]) -> List[dict]:
        groups = {}
        for m in mappings:
            if m.get("mapping_type") == "B" and m.get("group_first"):
                folder = m.get("shared_group", "")
                group_slides = [
                    m2["slide_num"]
                    for m2 in mappings
                    if m2.get("mapping_type") == "B" and m2.get("shared_group") == folder
                ]
                groups[folder] = {"folder": folder, "slides": group_slides, "slide_count": len(group_slides)}
        return list(groups.values())

    def _normalize_name(self, name: str) -> str:
        import unicodedata
        name = unicodedata.normalize("NFKC", name)
        return re.sub(r"\s+", " ", name).strip().lower()

    def _normalize_activity_id(self, text: str) -> str:
        text = re.sub(r"\s+", " ", text).strip()
        text = re.sub(
            r"(TS_FCT_USS_\d+)\s+Action\s*-?\s*(\d+)",
            r"\1_Action\2", text, flags=re.IGNORECASE,
        )
        return re.sub(r"\s+", "_", text)

    def _parse_driver_activity_line(self, da_text: str) -> List[str]:
        da_text = re.sub(r"Driver\s*Activity\s*:\s*", "", da_text, flags=re.IGNORECASE).strip()
        base_matches = list(re.finditer(r"TS_FCT_USS_\d+", da_text, re.IGNORECASE))
        if not base_matches:
            return []

        results = []
        if len(base_matches) == 1:
            base = base_matches[0].group(0)
            after_base = da_text[base_matches[0].end():]
            action_matches = re.findall(r"Action\s*-?\s*(\d+)", after_base, re.IGNORECASE)
            if action_matches:
                for an in action_matches:
                    results.append(f"{base}_Action{an}")
            else:
                results.append(base)
        else:
            parts = re.split(r"\s*&\s*", da_text)
            for part in parts:
                normalized = self._normalize_activity_id(part.strip())
                if normalized:
                    results.append(normalized)
        return results

    def _count_images(self, folder_path: str) -> int:
        if not os.path.isdir(folder_path):
            return 0
        return sum(1 for f in os.listdir(folder_path) if f.lower().endswith((".png", ".jpg", ".jpeg")))

    def _fuzzy_match_folder(self, target: str, candidates: List[str]) -> Optional[str]:
        tl = target.lower().replace("_", "").replace(" ", "")
        for c in candidates:
            if tl == c.lower().replace("_", "").replace(" ", ""):
                return c
        return None

    # 图片最大尺寸（超出则等比缩放），PPT 不需要 4K 分辨率
    @property
    def max_image_dim(self) -> int:
        return self._cfg("image_processing", "max_dimension", default=1920)

    def _convert_image_if_needed(self, img_path: str) -> str:
        """MPO 等非标准格式转为 JPEG，限制分辨率以控制文件大小。"""
        pptx_supported_formats = {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}
        try:
            img = Image.open(img_path)
            fmt = img.format
            if fmt and fmt.upper() in pptx_supported_formats:
                img.close()
                return img_path
            # MPO 等非标准格式 → JPEG
            converted = self._pil_to_jpeg(img)
            img.close()
            return converted
        except Exception:
            return img_path

    def _pil_to_jpeg(self, img: Image.Image) -> str:
        """PIL Image → 临时 JPEG，转换 RGBA/P 模式并限制分辨率。"""
        if img.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg
        elif img.mode == "P":
            img = img.convert("RGBA")
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        # 限制最大尺寸
        w, h = img.size
        ratio = min(self.max_image_dim / w, self.max_image_dim / h, 1.0)
        if ratio < 1:
            img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
        tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
        img.save(tmp.name, "JPEG", quality=85)
        tmp.close()
        return tmp.name
